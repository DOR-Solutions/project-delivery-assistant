"""Render a mitigation plan to a PowerPoint deck in the Pilz house format."""
import io

# Heathrow / MAX palette
TEAL = (0x00, 0x9E, 0x84)
INK = (0x10, 0x28, 0x3B)
RED = (0xD4, 0x37, 0x4C)
AMBER = (0xB0, 0x72, 0x0A)
GREY = (0x6B, 0x80, 0x93)
WHITE = (0xFF, 0xFF, 0xFF)
LIGHT = (0xEA, 0xF1, 0xF6)


def _pri_color(p: str):
    return {"critical": RED, "high": AMBER, "medium": (0xC9, 0xA2, 0x27), "low": TEAL}.get((p or "").lower(), GREY)


def build_pptx(plan: dict) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    SW, SH = prs.slide_width, prs.slide_height

    def rgb(c):
        return RGBColor(*c)

    def add_slide():
        return prs.slides.add_slide(blank)

    def title_bar(slide, text, sub=""):
        bar = slide.shapes.add_shape(1, 0, 0, SW, Inches(0.95))
        bar.fill.solid(); bar.fill.fore_color.rgb = rgb(INK); bar.line.fill.background()
        tf = bar.text_frame; tf.margin_left = Inches(0.4); tf.word_wrap = True
        p = tf.paragraphs[0]; r = p.add_run(); r.text = text
        r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = rgb(WHITE)
        if sub:
            p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = sub
            r2.font.size = Pt(12); r2.font.color.rgb = rgb(RGBColor(0xAE, 0xC8, 0xD8))

    def textbox(slide, text, top, size=14, color=INK, bold=False, left=Inches(0.5), width=None, height=Inches(1.0)):
        box = slide.shapes.add_textbox(left, top, width or (SW - Inches(1.0)), height)
        tf = box.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = rgb(color)
        return box

    def add_table(slide, headers, rows, top, col_widths, height=Inches(5.0)):
        nrows = len(rows) + 1
        ncols = len(headers)
        total_w = sum(col_widths)
        left = int((SW - total_w) / 2)
        gtbl = slide.shapes.add_table(nrows, ncols, left, top, total_w, height).table
        for j, w in enumerate(col_widths):
            gtbl.columns[j].width = w
        for j, h in enumerate(headers):
            c = gtbl.cell(0, j); c.text = h
            pf = c.text_frame.paragraphs[0]; pf.runs[0].font.size = Pt(11); pf.runs[0].font.bold = True
            pf.runs[0].font.color.rgb = rgb(WHITE); c.fill.solid(); c.fill.fore_color.rgb = rgb(TEAL)
        for i, row in enumerate(rows, 1):
            for j, val in enumerate(row):
                c = gtbl.cell(i, j); c.text = str(val)
                pf = c.text_frame.paragraphs[0]
                if pf.runs:
                    pf.runs[0].font.size = Pt(10); pf.runs[0].font.color.rgb = rgb(INK)
                c.fill.solid(); c.fill.fore_color.rgb = rgb(WHITE if i % 2 else LIGHT)
        return gtbl

    project = plan.get("project", "Mitigation Plan")

    # 1. Title
    s = add_slide()
    bg = s.shapes.add_shape(1, 0, 0, SW, SH); bg.fill.solid(); bg.fill.fore_color.rgb = rgb(INK); bg.line.fill.background()
    textbox(s, "MAX.ai — Mitigation Plan", Inches(2.6), size=40, color=WHITE, bold=True, left=Inches(0.8), width=SW - Inches(1.6))
    textbox(s, project, Inches(3.7), size=22, color=(0x6F, 0xE0, 0xC8), left=Inches(0.8), width=SW - Inches(1.6))
    textbox(s, "Terminal 5 Operations · UMP Project Intelligence", Inches(4.5), size=14, color=(0xAE, 0xC8, 0xD8), left=Inches(0.8))

    # 2. Objective
    s = add_slide(); title_bar(s, "Objective")
    textbox(s, plan.get("objective", ""), Inches(1.4), size=16, height=Inches(2.0))
    if plan.get("narrative"):
        textbox(s, plan["narrative"], Inches(3.4), size=13, color=GREY, height=Inches(3.0))

    # 3. Mitigation actions
    acts = plan.get("mitigation_actions", [])
    if acts:
        s = add_slide(); title_bar(s, "Detailed Mitigation Strategy")
        rows = [[a.get("area", ""), a.get("action", ""), a.get("mitigation", ""), a.get("responsibility", ""), (a.get("priority", "") or "").upper()] for a in acts[:12]]
        add_table(s, ["Area", "Action", "Mitigation", "Responsibility", "Pri"],
                  rows, Inches(1.2), [Inches(1.4), Inches(2.6), Inches(5.0), Inches(2.0), Inches(1.0)])

    # 4. FMEA
    fmea = plan.get("fmea", [])
    if fmea:
        s = add_slide(); title_bar(s, "Failure Modes & Effects Analysis (FMEA)")
        rows = [[f.get("process", ""), f.get("failure_mode", ""), f.get("effect", ""), str(f.get("severity", "")), f.get("controls", "")] for f in fmea[:10]]
        add_table(s, ["Process", "Failure Mode", "Effect", "Sev", "Controls"],
                  rows, Inches(1.2), [Inches(1.4), Inches(3.2), Inches(2.8), Inches(0.8), Inches(3.8)])

    # 5. Access windows
    aw = plan.get("access_windows", [])
    if aw:
        s = add_slide(); title_bar(s, "Access Windows & Schedule Impact")
        rows = [[str(a.get("item", "")), a.get("area", ""), a.get("access", ""), a.get("start", ""), a.get("finish", ""), a.get("original_duration", ""), a.get("new_duration", "")] for a in aw]
        add_table(s, ["Item", "Line/Area", "Access", "Start", "Finish", "Original", "New"],
                  rows, Inches(1.3), [Inches(0.8), Inches(2.0), Inches(3.0), Inches(1.6), Inches(1.6), Inches(1.3), Inches(1.3)])

    # 6. Predicted risks + to-do
    pr = plan.get("predicted_risks", [])
    if pr:
        s = add_slide(); title_bar(s, "Predicted Risks")
        rows = [[(r.get("band", "") or "").upper(), r.get("title", ""), r.get("rationale", "")] for r in pr[:10]]
        add_table(s, ["Band", "Risk", "Rationale"], rows, Inches(1.2), [Inches(1.4), Inches(3.6), Inches(6.0)])

    todo = plan.get("todo", [])
    if todo:
        s = add_slide(); title_bar(s, "Project-Manager To-Do")
        rows = [[(t.get("priority", t.get("pri", "")) or "").upper(), t.get("text", ""), t.get("owner", "")] for t in todo[:12]]
        add_table(s, ["Pri", "Action", "Owner"], rows, Inches(1.2), [Inches(1.2), Inches(7.4), Inches(2.4)])

    # 7. Command & control + contingency
    s = add_slide(); title_bar(s, "Command & Control · Contingency")
    textbox(s, "Command & Control", Inches(1.3), size=16, bold=True, color=TEAL)
    textbox(s, plan.get("command_control", ""), Inches(1.9), size=13, height=Inches(1.6))
    textbox(s, "Contingency", Inches(3.7), size=16, bold=True, color=TEAL)
    textbox(s, plan.get("contingency", ""), Inches(4.3), size=13, height=Inches(1.6))

    # 8. Approvals
    appr = plan.get("approvals", [])
    if appr:
        s = add_slide(); title_bar(s, "Stakeholder Approvals")
        rows = [[a.get("name", ""), a.get("company", ""), a.get("role", ""), "", ""] for a in appr]
        add_table(s, ["Name", "Company", "Role", "Signature", "Date"],
                  rows, Inches(1.2), [Inches(2.6), Inches(2.2), Inches(3.2), Inches(2.0), Inches(1.6)])

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
