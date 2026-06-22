"""Server-side file parsing: PDF, XLSX, DOCX, CSV/TXT -> plain text."""
import io
from typing import Tuple

def parse_pdf(data: bytes) -> str:
    import pdfplumber
    out = []
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for page in pdf.pages[:50]:
            out.append(page.extract_text() or "")
    return "\n\n".join(out).strip()

def parse_xlsx(data: bytes) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True)
    chunks = []
    for ws in wb.worksheets:
        chunks.append(f"=== Sheet: {ws.title} ===")
        for row in ws.iter_rows(values_only=True):
            vals = [str(v) for v in row if v is not None]
            if vals:
                chunks.append(", ".join(vals))
    return "\n".join(chunks).strip()

def parse_docx(data: bytes) -> str:
    import docx
    d = docx.Document(io.BytesIO(data))
    parts = [p.text for p in d.paragraphs if p.text.strip()]
    for t in d.tables:
        for r in t.rows:
            parts.append(" | ".join(c.text.strip() for c in r.cells))
    return "\n".join(parts).strip()

def parse_file(filename: str, data: bytes) -> Tuple[str, str]:
    """Returns (text, kind)."""
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    if ext == "pdf":
        return parse_pdf(data), "pdf"
    if ext in ("xlsx", "xlsm", "xls"):
        return parse_xlsx(data), "xlsx"
    if ext in ("docx", "doc"):
        return parse_docx(data), "docx"
    if ext in ("pptx", "ppt"):
        return parse_pptx(data), "pptx"
    if ext in ("vtt", "srt"):
        return parse_vtt(data), "transcript"
    return data.decode("utf-8", errors="replace"), ("csv" if ext in ("csv", "tsv") else "text")


def parse_vtt(data: bytes) -> str:
    """Microsoft Teams / WebVTT (and SRT) transcript -> clean speaker-tagged text.

    Strips the WEBVTT header, cue indices and timestamp lines, unwraps the
    ``<v Speaker>…</v>`` voice tags Teams emits, and merges consecutive turns by
    the same speaker so the result reads like minutes (``Speaker: text``)."""
    import re
    raw = data.decode("utf-8-sig", errors="replace")
    turns: list[list] = []  # [speaker_or_None, text]
    for ln in raw.splitlines():
        s = ln.strip()
        if not s or s.upper() == "WEBVTT" or s.startswith("NOTE") or "-->" in s or re.match(r"^\d+$", s):
            continue
        m = re.search(r"<v\s+([^>]+)>(.*)", s)
        if m:
            sp, txt = m.group(1).strip(), m.group(2)
        else:
            mm = re.match(r"^([A-Z][\w'.\-]+(?:\s+[A-Z][\w'.\-]+){0,3}):\s*(.*)$", s)
            sp, txt = (mm.group(1).strip(), mm.group(2)) if mm else (None, s)
        txt = re.sub(r"<[^>]+>", "", txt).strip()
        if not txt:
            continue
        if turns and turns[-1][0] == sp:
            turns[-1][1] += " " + txt
        else:
            turns.append([sp, txt])
    return "\n".join((f"{sp}: {txt}" if sp else txt) for sp, txt in turns).strip()


def parse_pptx(data: bytes) -> str:
    """Slide text + tables from PowerPoint. Uses python-pptx when available,
    otherwise falls back to reading the slide XML directly (so it works even
    before the dependency is installed)."""
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        chunks = []
        for i, slide in enumerate(prs.slides, 1):
            chunks.append(f"=== Slide {i} ===")
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    chunks.append(shape.text_frame.text.strip())
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [c.text.strip() for c in row.cells]
                        if any(cells):
                            chunks.append(" | ".join(cells))
        return "\n".join(chunks).strip()
    except Exception:
        # XML fallback: pull every <a:t> text run, slide by slide
        import zipfile, re, html
        out = []
        with zipfile.ZipFile(io.BytesIO(data)) as z:
            slides = sorted([n for n in z.namelist() if re.match(r"ppt/slides/slide\d+\.xml$", n)],
                            key=lambda n: int(re.search(r"(\d+)", n).group()))
            for i, sn in enumerate(slides, 1):
                runs = re.findall(r"<a:t>(.*?)</a:t>", z.read(sn).decode("utf-8", "replace"), re.S)
                text = html.unescape(" ".join(r.strip() for r in runs if r.strip()))
                if text:
                    out.append(f"=== Slide {i} ===\n{text}")
        return "\n".join(out).strip()
